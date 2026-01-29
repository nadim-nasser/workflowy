#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Define colors
DARK_BG = RGBColor(15, 23, 42)  # slate-900
BLUE = RGBColor(59, 130, 246)
PURPLE = RGBColor(139, 92, 246)
EMERALD = RGBColor(16, 185, 129)
AMBER = RGBColor(245, 158, 11)
WHITE = RGBColor(255, 255, 255)
SLATE_300 = RGBColor(203, 213, 225)
SLATE_400 = RGBColor(148, 163, 184)
DELOITTE_GREEN = RGBColor(134, 188, 37)

def add_dark_slide(prs, title, subtitle=None, content=None, accent_color=BLUE):
    """Add a dark-themed slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Dark background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()

    return slide

def add_light_slide(prs):
    """Add a light-themed slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Light background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(248, 250, 252)  # slate-50
    background.line.fill.background()

    return slide

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    """Add a text box to a slide"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox

# Slide 1: Cover
slide = add_dark_slide(prs, "Cover")
add_text_box(slide, 1, 1.5, 11, 0.5, "Workflowy × Deloitte", font_size=24, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 1, 2.5, 11, 1, "Executive AI Training", font_size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.8, 11, 0.5, "Half-Day Intensive for Deloitte Leaders", font_size=24, color=SLATE_300, align=PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.5, 11, 0.5, "January 2026", font_size=18, color=SLATE_400, align=PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.5, 11, 0.5, "Trusted by leaders at Google, KPMG, TD Bank, RBC", font_size=14, color=SLATE_400, align=PP_ALIGN.CENTER)

# Slide 2: The Challenge
slide = add_light_slide(prs)
add_text_box(slide, 0.8, 0.5, 3, 0.4, "THE PROBLEM", font_size=12, bold=True, color=BLUE)
add_text_box(slide, 0.8, 0.9, 11, 1, "AI capability is compounding.\nEnterprise readiness is not.", font_size=40, bold=True, color=RGBColor(15, 23, 42))
add_text_box(slide, 0.8, 2.2, 10, 0.5, "Most deployments stall at chat interfaces, copilots, or fragile pilots that collapse under real operational pressure.", font_size=16, color=RGBColor(71, 85, 105))

# Three problem cards
problems = [
    ("Pilots Everywhere, Impact Nowhere", "AI lives at the edges: chat interfaces and brittle demos."),
    ("Chat ≠ Capability", "Teams ask AI questions. They don't build with it."),
    ("Intelligence Without Infrastructure", "Powerful AI, no systems to deploy it.")
]
for i, (title, desc) in enumerate(problems):
    x = 0.8 + i * 4
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3), Inches(3.7), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    add_text_box(slide, x + 0.2, 3.2, 3.3, 0.4, title, font_size=14, bold=True, color=RGBColor(15, 23, 42))
    add_text_box(slide, x + 0.2, 3.6, 3.3, 0.8, desc, font_size=12, color=RGBColor(71, 85, 105))

# Stats
stats = [("72%", "report AI skills gaps as top barrier"), ("$4.4T", "annual GenAI value, mostly untapped"), ("3x", "gains with structured training")]
for i, (num, desc) in enumerate(stats):
    x = 1.5 + i * 4
    add_text_box(slide, x, 5, 3, 0.5, num, font_size=32, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text_box(slide, x, 5.6, 3, 0.5, desc, font_size=11, color=RGBColor(71, 85, 105), align=PP_ALIGN.CENTER)

# Slide 3: Our Solution
slide = add_dark_slide(prs, "Solution")
add_text_box(slide, 0.8, 0.5, 3, 0.4, "OUR SOLUTION", font_size=12, bold=True, color=BLUE)
add_text_box(slide, 0.8, 0.9, 11, 1.2, "We don't teach AI.\nWe extract and operationalize\nyour workflows.", font_size=36, bold=True, color=WHITE)
add_text_box(slide, 0.8, 2.5, 10, 0.8, "The real value lives in the tacit knowledge of your people: the approvals, the edge cases, the \"how we actually do things here.\"", font_size=16, color=SLATE_300)

solutions = [
    ("Workflow Extraction", "Surface processes and institutional knowledge, systematize them for AI."),
    ("Live Builds", "Build functional AI solutions together, deployed in hours, not months."),
    ("Your Platforms", "Master Microsoft Copilot, Google Gemini, AWS with Claude. No new vendors.")
]
for i, (title, desc) in enumerate(solutions):
    x = 0.8 + i * 4
    add_text_box(slide, x, 3.5, 3.5, 0.4, title, font_size=18, bold=True, color=WHITE)
    add_text_box(slide, x, 4, 3.5, 0.8, desc, font_size=13, color=SLATE_400)

# Slide 4: The Evolution
slide = add_dark_slide(prs, "Evolution")
add_text_box(slide, 0.8, 0.5, 3, 0.4, "THE EVOLUTION", font_size=12, bold=True, color=BLUE)
add_text_box(slide, 0.8, 0.9, 11, 0.8, "From chatbots to autonomous agents", font_size=40, bold=True, color=WHITE)
add_text_box(slide, 0.8, 1.8, 10, 0.5, "Leading enterprises are moving beyond chat. Can your people keep up?", font_size=16, color=SLATE_400)

stages = [
    ("Stage 1", "Chat Interface", "Human at keyboard. Reactive Q&A.", "Where most are today", SLATE_400),
    ("Stage 2", "Human Guided Agents", "AI completes multi-step tasks.", "Where leaders need to be", BLUE),
    ("Stage 3", "Autonomous Agents", "AI acts independently.", "Where enterprises are headed", PURPLE)
]
for i, (stage, title, desc, note, color) in enumerate(stages):
    x = 0.8 + i * 4
    add_text_box(slide, x, 2.8, 3.5, 0.3, stage, font_size=11, color=color)
    add_text_box(slide, x, 3.1, 3.5, 0.4, title, font_size=18, bold=True, color=WHITE)
    add_text_box(slide, x, 3.5, 3.5, 0.5, desc, font_size=13, color=SLATE_400)
    add_text_box(slide, x, 4.1, 3.5, 0.3, note, font_size=11, color=color)

add_text_box(slide, 0.8, 5.5, 11.5, 0.8, "Your clients are already building agentic platforms. Major banks are deploying agent architectures with LLM gateways.", font_size=14, color=SLATE_300, align=PP_ALIGN.CENTER)

# Slide 5: Skills Gap
slide = add_light_slide(prs)
add_text_box(slide, 0.8, 0.5, 3, 0.4, "THE SKILLS GAP", font_size=12, bold=True, color=BLUE)
add_text_box(slide, 0.8, 0.9, 11, 0.8, "Three levels of AI capability", font_size=40, bold=True, color=RGBColor(15, 23, 42))
add_text_box(slide, 0.8, 1.8, 10, 0.5, "Enterprises need people across all three levels. Most have only Level 1.", font_size=16, color=RGBColor(71, 85, 105))

levels = [
    ("No Code", "AI User", "Everyday use, boosting productivity.", ["Summarize meeting notes", "Draft proposals", "Research benchmarks"], BLUE),
    ("Low Code", "AI Builder", "Power users driving automation.", ["Build client analyzer agents", "Automate workflows", "Create pipelines"], PURPLE),
    ("Pro Code", "AI Architect", "Advanced programmatic control.", ["Audit automation with HITL", "Native agents for delivery", "Enterprise integration"], EMERALD)
]
for i, (badge, title, desc, items, color) in enumerate(levels):
    x = 0.8 + i * 4
    add_text_box(slide, x, 2.5, 1.2, 0.3, badge, font_size=11, bold=True, color=WHITE)
    add_text_box(slide, x, 2.9, 3.5, 0.4, title, font_size=20, bold=True, color=RGBColor(15, 23, 42))
    add_text_box(slide, x, 3.4, 3.5, 0.4, desc, font_size=12, color=RGBColor(71, 85, 105))
    for j, item in enumerate(items):
        add_text_box(slide, x, 3.9 + j * 0.35, 3.5, 0.3, f"• {item}", font_size=11, color=RGBColor(71, 85, 105))

add_text_box(slide, 0.8, 5.5, 11.5, 0.8, "The risk isn't AI replacing your consultants.\nIt's competitors whose consultants know how to use AI winning your engagements.", font_size=14, bold=True, color=RGBColor(15, 23, 42), align=PP_ALIGN.CENTER)

# Slide 6: Our Approach
slide = add_light_slide(prs)
add_text_box(slide, 0.8, 0.5, 3, 0.4, "OUR APPROACH", font_size=12, bold=True, color=BLUE)
add_text_box(slide, 0.8, 0.9, 11, 0.8, "From awareness to operational capability", font_size=40, bold=True, color=RGBColor(15, 23, 42))

# Example functions
add_text_box(slide, 0.8, 1.8, 6, 0.4, "Example Functions You'll Build:", font_size=14, bold=True, color=RGBColor(71, 85, 105))
functions = ["Client Brief Generator", "Research Synthesizer", "Data Story Builder", "Meeting Prep Agent"]
for i, func in enumerate(functions):
    x = 0.8 + i * 3
    add_text_box(slide, x, 2.2, 2.8, 0.3, func, font_size=12, color=RGBColor(15, 23, 42))

steps = [
    ("1", "Assess & Baseline", ["Pre-session readiness assessment", "Individual skill baselining", "Gap & opportunity mapping"], BLUE),
    ("2", "Custom Training", ["Role specific content", "Your approved tools", "Hands on building"], PURPLE),
    ("3", "Measure & Scale", ["Post-session skill assessment", "Champion identification", "ROI reporting & scale plan"], EMERALD)
]
for i, (num, title, items, color) in enumerate(steps):
    x = 0.8 + i * 4
    add_text_box(slide, x, 3, 0.4, 0.4, num, font_size=18, bold=True, color=WHITE)
    add_text_box(slide, x + 0.5, 3, 3, 0.4, title, font_size=18, bold=True, color=RGBColor(15, 23, 42))
    for j, item in enumerate(items):
        add_text_box(slide, x, 3.5 + j * 0.35, 3.5, 0.3, f"• {item}", font_size=11, color=RGBColor(71, 85, 105))

add_text_box(slide, 0.8, 5.2, 11.5, 0.4, "Training on Your Tools: Claude • Gemini • Copilot • Gen D • Cursor", font_size=14, color=SLATE_400, align=PP_ALIGN.CENTER)

# Slide 7: Workshop Details
slide = add_light_slide(prs)
add_text_box(slide, 0.8, 0.5, 3, 0.4, "PROPOSED WORKSHOP", font_size=12, bold=True, color=BLUE)
add_text_box(slide, 0.8, 0.9, 11, 0.8, "Half-Day Executive AI Intensive", font_size=40, bold=True, color=RGBColor(15, 23, 42))
add_text_box(slide, 0.8, 1.8, 10, 0.5, "A focused, hands-on session designed to shift your leaders from curious to capable.", font_size=16, color=RGBColor(71, 85, 105))

details = [
    ("20 Participants", "Service line and industry leaders"),
    ("Pre-Session Discovery", "Tools, constraints, priorities"),
    ("Custom Curriculum", "Your approved tech stack"),
    ("Live Builds", "Create something real")
]
for i, (title, desc) in enumerate(details):
    x = 0.8 + i * 3
    add_text_box(slide, x, 2.8, 2.8, 0.4, title, font_size=16, bold=True, color=RGBColor(15, 23, 42))
    add_text_box(slide, x, 3.2, 2.8, 0.4, desc, font_size=12, color=RGBColor(71, 85, 105))

add_text_box(slide, 0.8, 4.5, 11.5, 0.8, "WorkflowyOS Learning Platform Included:\nInteractive courses, progress tracking, AI coaching assistant, and post-session resources.", font_size=14, color=RGBColor(15, 23, 42), align=PP_ALIGN.CENTER)

# Slide 8: Program Curriculum
slide = add_dark_slide(prs, "Curriculum")
add_text_box(slide, 0.8, 0.5, 3, 0.4, "PROGRAM CURRICULUM", font_size=12, bold=True, color=BLUE)
add_text_box(slide, 0.8, 0.9, 11, 0.8, "Four hours, four modules", font_size=40, bold=True, color=WHITE)
add_text_box(slide, 0.8, 1.7, 10, 0.4, "50% interactive discussion, 50% hands-on application.", font_size=14, color=SLATE_400)

modules = [
    ("1", "AI Landscape & Reality Check", "45 min", "Cut through the hype. What AI can actually do today.", BLUE),
    ("2", "AI Fluency Fundamentals", "60 min", "Master prompt engineering with RACE framework.", PURPLE),
    ("3", "Hands-On Building", "90 min", "Live build exercises using your approved tools.", EMERALD),
    ("4", "Workflow Integration", "45 min", "Leave with a personalized AI adoption roadmap.", AMBER)
]
for i, (num, title, time, desc, color) in enumerate(modules):
    x = 0.8 + (i % 2) * 6
    y = 2.3 + (i // 2) * 1.8
    add_text_box(slide, x, y, 0.4, 0.4, num, font_size=18, bold=True, color=color)
    add_text_box(slide, x + 0.5, y, 4.5, 0.4, title, font_size=16, bold=True, color=WHITE)
    add_text_box(slide, x + 5.2, y, 0.8, 0.3, time, font_size=10, color=SLATE_400)
    add_text_box(slide, x + 0.5, y + 0.4, 5, 0.4, desc, font_size=12, color=SLATE_400)

# Slide 9: Outcomes
slide = add_light_slide(prs)
add_text_box(slide, 0.8, 0.5, 3, 0.4, "PARTICIPANT OUTCOMES", font_size=12, bold=True, color=BLUE)
add_text_box(slide, 0.8, 0.9, 11, 0.8, "What your leaders leave with", font_size=40, bold=True, color=RGBColor(15, 23, 42))

outcomes = [
    ("Mindset Shift", ['From "AI as search" to "AI as assistant"', "Clear AI capabilities vs. limitations", "Confidence with clients and teams"], BLUE),
    ("Practical Skills", ["RACE framework for prompts", "Identify 3+ automation opportunities", "Navigate security boundaries"], PURPLE),
    ("Tangible Deliverables", ["Pre/post skill progression report", "Working AI solution built in session", "Role-specific prompt library"], EMERALD),
    ("Post-Session Resources", ["Curated resource kit & templates", "30-day email follow-up series", "Partner network access"], AMBER)
]
for i, (title, items, color) in enumerate(outcomes):
    x = 0.8 + (i % 2) * 6
    y = 1.8 + (i // 2) * 2.2
    add_text_box(slide, x, y, 5.5, 0.4, title, font_size=16, bold=True, color=RGBColor(15, 23, 42))
    for j, item in enumerate(items):
        add_text_box(slide, x, y + 0.4 + j * 0.35, 5.5, 0.3, f"✓ {item}", font_size=11, color=RGBColor(71, 85, 105))

# Slide 10: Measurement & ROI
slide = add_dark_slide(prs, "Measurement")
add_text_box(slide, 0.8, 0.5, 3, 0.4, "MEASURABLE IMPACT", font_size=12, bold=True, color=EMERALD)
add_text_box(slide, 0.8, 0.9, 11, 0.8, "Skills-based assessment. Proven ROI.", font_size=40, bold=True, color=WHITE)
add_text_box(slide, 0.8, 1.8, 10, 0.5, "We measure skill progression before and after, giving you concrete data on your AI readiness investment.", font_size=14, color=SLATE_300)

add_text_box(slide, 0.8, 2.6, 5.5, 0.4, "Pre/Post AI Readiness Assessment", font_size=16, bold=True, color=WHITE)
assessment_items = ["Baseline assessment before training", "Post-session competency evaluation", "30/60/90-day follow-up measurements"]
for i, item in enumerate(assessment_items):
    add_text_box(slide, 0.8, 3 + i * 0.35, 5.5, 0.3, f"✓ {item}", font_size=12, color=SLATE_300)

add_text_box(slide, 7, 2.6, 5.5, 0.4, "Skills-Based Talent Reporting", font_size=16, bold=True, color=WHITE)
talent_items = ["Identify high-potential AI adopters", "Skills heat map across roles", "Targeted development paths"]
for i, item in enumerate(talent_items):
    add_text_box(slide, 7, 3 + i * 0.35, 5.5, 0.3, f"✓ {item}", font_size=12, color=SLATE_300)

add_text_box(slide, 0.8, 4.5, 11.5, 0.4, "The Workflowy Difference: Live Builds + Measured Impact", font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.8, 5, 11.5, 0.4, "Other programs measure awareness. We measure capability because participants build working solutions.", font_size=14, color=SLATE_300, align=PP_ALIGN.CENTER)

metrics = [("100%", "Build something real"), ("Pre+Post", "Skill progression"), ("Champions", "For scale-out")]
for i, (num, desc) in enumerate(metrics):
    x = 2 + i * 3.5
    add_text_box(slide, x, 5.7, 3, 0.5, num, font_size=28, bold=True, color=EMERALD, align=PP_ALIGN.CENTER)
    add_text_box(slide, x, 6.2, 3, 0.3, desc, font_size=11, color=SLATE_400, align=PP_ALIGN.CENTER)

# Slide 11: Why Workflowy
slide = add_light_slide(prs)
add_text_box(slide, 0.8, 0.5, 3, 0.4, "WHY WORKFLOWY", font_size=12, bold=True, color=BLUE)
add_text_box(slide, 0.8, 0.9, 11, 0.8, "What makes us different", font_size=40, bold=True, color=RGBColor(15, 23, 42))

differentiators = [
    ("Hands-On, Not Theory", "Participants build real tools during the session. No death by PowerPoint.", BLUE),
    ("Your Tools, Your Constraints", "We work within your approved tech stack and compliance requirements.", PURPLE),
    ("Results, Not Inspiration", "Measurable skill progression. Clear ROI. Monday morning behavior change.", EMERALD)
]
for i, (title, desc, color) in enumerate(differentiators):
    x = 0.8 + i * 4
    add_text_box(slide, x, 2, 3.5, 0.4, title, font_size=18, bold=True, color=RGBColor(15, 23, 42))
    add_text_box(slide, x, 2.5, 3.5, 0.8, desc, font_size=13, color=RGBColor(71, 85, 105))

add_text_box(slide, 0.8, 4.5, 11.5, 0.8, '"Most AI training stops at awareness. We stop at execution."\n\nWe don\'t just teach people about AI. We ensure they use it.', font_size=16, color=RGBColor(15, 23, 42), align=PP_ALIGN.CENTER)

# Slide 12: Team
slide = add_light_slide(prs)
add_text_box(slide, 0.8, 0.5, 3, 0.4, "OUR TEAM", font_size=12, bold=True, color=BLUE)
add_text_box(slide, 0.8, 0.9, 11, 0.8, "Enterprise experience meets AI expertise", font_size=40, bold=True, color=RGBColor(15, 23, 42))

team = [
    ("Nadim Nasser", "CEO & Head of Training", ["15+ years education & tech", "Former Head of Ed at Prequel", "Trained 500+ on AI"]),
    ("Drew Baillie", "Senior AI Consultant", ["25+ years AI transformation", "Former KPMG AI lead", "Board Director"]),
    ("Azim Ahmed", "VP Engineering", ["Head of Eng at Lazer", "5+ years leading teams", "AI implementation"])
]
for i, (name, role, items) in enumerate(team):
    x = 0.8 + i * 4
    add_text_box(slide, x, 1.9, 3.5, 0.4, name, font_size=18, bold=True, color=RGBColor(15, 23, 42), align=PP_ALIGN.CENTER)
    add_text_box(slide, x, 2.3, 3.5, 0.3, role, font_size=12, color=BLUE, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text_box(slide, x, 2.7 + j * 0.3, 3.5, 0.3, f"• {item}", font_size=11, color=RGBColor(71, 85, 105))

add_text_box(slide, 0.8, 4.2, 11.5, 0.4, "Advisory Board", font_size=16, bold=True, color=RGBColor(15, 23, 42), align=PP_ALIGN.CENTER)

advisors = [
    ("Armughan Ahmad", "Executive Chairman", ["30 year enterprise career", "Advisory: OpenAI, Telus, ServiceNow"]),
    ("Arif Bhanji", "Co-Founder, Lazer Technologies", ["Former Monitor Deloitte Consultant", "Y Combinator alumni"])
]
for i, (name, role, items) in enumerate(advisors):
    x = 2.5 + i * 5
    add_text_box(slide, x, 4.7, 4, 0.4, name, font_size=16, bold=True, color=RGBColor(15, 23, 42), align=PP_ALIGN.CENTER)
    add_text_box(slide, x, 5.1, 4, 0.3, role, font_size=11, color=AMBER, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text_box(slide, x, 5.4 + j * 0.3, 4, 0.3, f"• {item}", font_size=10, color=RGBColor(71, 85, 105))

# Slide 13: Investment
slide = add_dark_slide(prs, "Investment")
add_text_box(slide, 0.8, 0.5, 3, 0.4, "INVESTMENT", font_size=12, bold=True, color=BLUE)
add_text_box(slide, 0.8, 1.2, 11, 0.8, "Half-Day Executive Intensive", font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, 0.8, 2.5, 11, 1, "$50,000", font_size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.8, 3.5, 11, 0.4, "For 20 participants ($2,500/executive)", font_size=16, color=SLATE_400, align=PP_ALIGN.CENTER)

includes = [
    ("Included", ["Pre-session discovery call", "Customized curriculum", "4-hour live intensive", "All materials & templates"]),
    ("Deliverables", ["Pre/post skill assessment", "AI readiness report", "Resource kit access", "30-day follow-up"]),
    ("Ongoing Support", ["WorkflowyOS platform access", "Email support for 30 days", "Champion certification", "Scale planning session"])
]
for i, (title, items) in enumerate(includes):
    x = 0.8 + i * 4
    add_text_box(slide, x, 4.2, 3.5, 0.4, title, font_size=14, bold=True, color=WHITE)
    for j, item in enumerate(items):
        add_text_box(slide, x, 4.6 + j * 0.35, 3.5, 0.3, f"• {item}", font_size=11, color=SLATE_300)

# Slide 14: Next Steps / Closing
slide = add_dark_slide(prs, "Next Steps")
add_text_box(slide, 0.8, 0.5, 3, 0.4, "NEXT STEPS", font_size=12, bold=True, color=BLUE)
add_text_box(slide, 0.8, 1.5, 11, 0.8, "Ready to transform your team's\nAI capabilities?", font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

next_steps = [
    ("1", "Schedule Discovery Call", "15-minute alignment on goals and constraints"),
    ("2", "Curriculum Customization", "We adapt training to your tools and use cases"),
    ("3", "Session Delivery", "Half-day intensive with your 20 leaders"),
    ("4", "Measure & Scale", "Assessment results and scale-out recommendations")
]
for i, (num, title, desc) in enumerate(next_steps):
    x = 0.8 + (i % 2) * 6
    y = 3.2 + (i // 2) * 1.2
    add_text_box(slide, x, y, 0.4, 0.4, num, font_size=18, bold=True, color=BLUE)
    add_text_box(slide, x + 0.5, y, 5, 0.4, title, font_size=16, bold=True, color=WHITE)
    add_text_box(slide, x + 0.5, y + 0.4, 5, 0.4, desc, font_size=12, color=SLATE_400)

add_text_box(slide, 0.8, 6, 11, 0.4, "nadim@workflowy.ai  |  workflowy.ai", font_size=18, color=SLATE_300, align=PP_ALIGN.CENTER)

# Save the presentation
prs.save('/Users/nadimnasser/GitHub/workflowy/Deloitte_AI_Training_Proposal.pptx')
print("PowerPoint created: Deloitte_AI_Training_Proposal.pptx")
